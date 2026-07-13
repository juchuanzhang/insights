from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG = RGBColor(0x0D, 0x1B, 0x2A)
TITLE_CLR = RGBColor(0xFF, 0xFF, 0xFF)
BODY_CLR = RGBColor(0xE0, 0xE0, 0xE0)
ACCENT = RGBColor(0x00, 0xBC, 0xD4)
WARN = RGBColor(0xFF, 0x6B, 0x6B)
OK = RGBColor(0x4C, 0xAF, 0x50)
SUBTITLE_CLR = RGBColor(0x90, 0xA4, 0xAE)
CODE_CLR = RGBColor(0xA5, 0xD6, 0xA7)
DIM = RGBColor(0x78, 0x90, 0x9C)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG

def add_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        shape.fill.fore_color.brightness = alpha
    return shape

def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)

def set_text(tf, text, size=18, color=BODY_CLR, bold=False, align=PP_ALIGN.LEFT):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return p

def add_para(tf, text, size=18, color=BODY_CLR, bold=False, align=PP_ALIGN.LEFT, space_before=Pt(6), space_after=Pt(2)):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    p.space_after = space_after
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return p

def title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT)
    tb = add_textbox(slide, Inches(1), Inches(2.2), Inches(11), Inches(2))
    set_text(tb.text_frame, title, size=36, color=TITLE_CLR, bold=True, align=PP_ALIGN.CENTER)
    if subtitle:
        add_para(tb.text_frame, subtitle, size=20, color=SUBTITLE_CLR, align=PP_ALIGN.CENTER)
    tb = add_textbox(slide, Inches(1), Inches(5), Inches(11), Inches(1))
    set_text(tb.text_frame, "Enter, Exit, Page Fault, Leak‡: Testing Isolation Boundaries for Microarchitectural Leaks", size=14, color=DIM, align=PP_ALIGN.CENTER)
    return slide

def section_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_rect(slide, Inches(4), Inches(2.8), Inches(5.333), Inches(0.06), ACCENT)
    tb = add_textbox(slide, Inches(1), Inches(2), Inches(11), Inches(1.5))
    set_text(tb.text_frame, title, size=30, color=TITLE_CLR, bold=True, align=PP_ALIGN.CENTER)
    if subtitle:
        add_para(tb.text_frame, subtitle, size=18, color=SUBTITLE_CLR, align=PP_ALIGN.CENTER)
    return slide

def content_slide(title, bullets=None, sub_title="", code_block=None, table=None, two_col=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_rect(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.06), ACCENT)
    tb = add_textbox(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.6))
    set_text(tb.text_frame, title, size=26, color=TITLE_CLR, bold=True)
    if sub_title:
        add_para(tb.text_frame, sub_title, size=16, color=ACCENT, bold=False)

    y = Inches(0.85)

    if two_col:
        left_items, right_items = two_col
        tb_l = add_textbox(slide, Inches(0.6), y, Inches(5.8), Inches(6))
        tb_r = add_textbox(slide, Inches(6.8), y, Inches(5.8), Inches(6))
        for item in left_items:
            if item.startswith("!"):
                add_para(tb_l.text_frame, item[1:], size=16, color=WARN, bold=True)
            elif item.startswith("+"):
                add_para(tb_l.text_frame, item[1:], size=16, color=OK, bold=True)
            else:
                add_para(tb_l.text_frame, item, size=16, color=BODY_CLR)
        for item in right_items:
            if item.startswith("!"):
                add_para(tb_r.text_frame, item[1:], size=16, color=WARN, bold=True)
            elif item.startswith("+"):
                add_para(tb_r.text_frame, item[1:], size=16, color=OK, bold=True)
            else:
                add_para(tb_r.text_frame, item, size=16, color=BODY_CLR)
        return slide

    if table:
        rows, cols = len(table), len(table[0])
        col_w = Inches(11.5 / cols)
        tbl_shape = slide.shapes.add_table(rows, cols, Inches(0.6), y, Inches(11.5), Inches(min(6, rows * 0.45)))
        tbl = tbl_shape.table
        for i, row_data in enumerate(table):
            for j, cell_text in enumerate(row_data):
                cell = tbl.cell(i, j)
                cell.text = ""
                p = cell.text_frame.paragraphs[0]
                run = p.add_run()
                run.text = cell_text
                if i == 0:
                    run.font.size = Pt(14)
                    run.font.color.rgb = ACCENT
                    run.font.bold = True
                else:
                    run.font.size = Pt(13)
                    run.font.color.rgb = BODY_CLR
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x15, 0x25, 0x35)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        return slide

    tb = add_textbox(slide, Inches(0.6), y, Inches(11.5), Inches(6))
    if bullets:
        for item in bullets:
            if item.startswith("!"):
                add_para(tb.text_frame, item[1:], size=16, color=WARN, bold=True)
            elif item.startswith("+"):
                add_para(tb.text_frame, item[1:], size=16, color=OK, bold=True)
            elif item.startswith(">>"):
                add_para(tb.text_frame, item[2:], size=15, color=CODE_CLR, bold=False, space_before=Pt(8))
            elif item.startswith("~"):
                add_para(tb.text_frame, item[1:], size=16, color=SUBTITLE_CLR, bold=True, space_before=Pt(10))
            else:
                add_para(tb.text_frame, item, size=16, color=BODY_CLR)

    if code_block:
        add_rect(slide, Inches(0.6), Inches(5.2), Inches(11.5), Inches(1.8), RGBColor(0x15, 0x25, 0x35))
        tb2 = add_textbox(slide, Inches(0.8), Inches(5.3), Inches(11), Inches(1.6))
        set_text(tb2.text_frame, code_block, size=13, color=CODE_CLR)

    return slide

# ============ SLIDES ============

title_slide(
    "Enter, Exit, Page Fault, Leak",
    "首个自动化跨域微架构泄漏测试框架"
)

section_slide("一、背景与动机")

content_slide("微架构泄漏：跨安全域的信息泄漏", bullets=[
    "~核心矛盾",
    "架构级隔离（页表、虚拟化、特权级）在逻辑层提供保护",
    "但共享微架构状态（缓存、CPU缓冲区、分支预测器）仍是泄漏通道",
    "~典型攻击",
    "Meltdown / Foreshadow / MDS / Spectre V2 / LVI ...",
    "~安全域（Security Domain）",
    "不同信任级别的执行环境：OS内核vs用户进程、hypervisor vs VM",
    "域转换：syscall/VM enter/exit等指令引起的上下文切换",
])

content_slide("防御碎片化与补丁屡次失效", bullets=[
    "~防御部署碎片化",
    "硬件层：部分由CPU微码补丁缓解（如Meltdown硬件修复）",
    "OS层：部分由软件缓解（KPTI、VERW缓冲区刷新）",
    "未修补：部分被认为安全影响低或缓解成本高而被忽略",
    "!补丁反复被证明不完整甚至完全无效",
    "Enhanced IBRS → BHI绕过",
    "FineIBT → BHI+IBPB绕过",
    "retpoline → RSB underflow绕过",
    "VERW → 特定条件下仍泄漏",
])

content_slide("现有工具无法跨域测试", sub_title="§2.3 现有方案局限",
    table=[
        ["工具", "域范围", "漏洞类型", "能否发现新类型", "能否跨域"],
        ["Revizor", "单域（内核）", "通用", "Yes", "No"],
        ["Scam-V", "单域（TrustZone）", "通用", "Yes", "No"],
        ["SpeechMiner", "单域", "Meltdown类定量", "No", "No"],
        ["Transynther", "双线程", "Meltdown/MDS变体", "No", "No"],
        ["RegCheck", "单域", "Meltdown-3a专用", "No", "No"],
        ["!核心差距：没有任何工具能系统性、自动化地测试跨域微架构隔离", "", "", "", ""],
    ])

section_slide("二、核心技术：Actor + MRT")

content_slide("Actor抽象：安全域的统一建模", sub_title="§3.1", bullets=[
    "~Actor定义",
    "Actor = 代码 + 私有数据 + 特权级 + CPU模式 + 内存权限配置",
    "~三种Actor模式（当前支持）",
    "Host-Kernel（特权级0，内核态）",
    "Host-User（特权级3，用户态）",
    "Guest-Kernel（VM内内核态）",
    "~域转换",
    "Actor之间通过预定义指令序列（syscall/VM enter/exit）转换",
    "每个test case包含一个main Actor + 一个或多个其他模式Actor",
])

content_slide("模板+宏机制：结构化与随机性的平衡", sub_title="§4.1",
    code_block=
"; 模板示例（双Actor：main + guest）\n"
".section .main\n"
".start:\n"
"  .macro.random_instructions.64:   ; 随机指令区\n"
"  VERW qword ptr [r14]             ; 缓解：缓冲区刷新\n"
"  .macro.set_h2g_target.vm_start:  ; 设置VM入口\n"
"  .macro.set_g2h_target.end:       ; 设置VM出口\n"
"  .macro.switch_h2g:               ; 域转换宏\n"
".end:\n"
"  .macro.fault_handler:            ; 异常处理\n\n"
".section .guest\n"
".vm_start:\n"
"  .macro.measurement_start:        ; 开始测量\n"
"  .macro.random_instructions.64:   ; 随机指令区\n"
"  .macro.measurement_end:          ; 结束测量\n"
"  .macro.switch_g2h:               ; 返回host",
    bullets=[
    "~模板（Template）",
    "汇编文件，定义Actor结构和域转换，不同Actor代码放在不同section",
    "~宏（Macro）",
    "伪指令（汇编中表现为NOP），在不同阶段差异化实例化",
    "执行器中：二进制补丁替换NOP为跳转到宏实现",
    "模型中：Unicorn回调拦截NOP地址调用对应实现",
    "~两大设计需求",
    "结构化：域转换/缓解需要固定指令序列（VMRESUME/VERW/L1D flush）",
    "统一化：执行器和模型需要不同实例化但统一描述",
])

content_slide("内存访问插桩", sub_title="§4.1.3", bullets=[
    "~新增插桩pass",
    "随机选择一个生成的内存访问，修改其地址指向其他Actor的内存",
    "专门针对Meltdown类泄漏：将用户态Actor的内存访问指向内核态Actor的内存",
    "~目的",
    "不用手写Meltdown gadget，工具自动探索'跨域非法访问'场景",
])

content_slide("MRT工作流程：五步循环", sub_title="§3.2", bullets=[
    "~模型关系测试（Model-Based Relational Testing）核心思路",
    "不直接比较硬件轨迹与模型预测，而是比较暴露的信息量",
    "~五步循环",
    "1. 程序生成 → 随机汇编指令序列（模板+配置约束）",
    "2. 输入生成 → 随机二进制数据初始化内存和寄存器",
    "3. 硬件执行 → 在目标CPU上执行，收集硬件轨迹（Flush+Reload）",
    "4. 模型执行 → 在Unicorn模拟器上执行，收集合同轨迹（预期泄漏）",
    "5. 泄漏比较 → 检查非干扰性质，过滤已知泄漏，报告未知泄漏",
    "~关键优势",
    "即使模型很简单，也能有效过滤已知泄漏并检测未知泄漏",
])

content_slide("合同式非干扰：形式化安全性质", sub_title="§6", bullets=[
    "~经典非干扰",
    "σ^1_A = σ^2_A ⇒ Measure(p, σ^1, μ) = Measure(p, σ^2, μ)",
    "受害者数据不应影响攻击者测量",
    "~合同式非干扰（实用化）",
    "σ^1_A = σ^2_A ∧ Contract(p, σ^V_1) = Contract(p, σ^V_2) ⇒ Measure(p, σ^1, μ) = Measure(p, σ^2, μ)",
    "~合同轨迹Contract(p, σ^V)",
    "模型预测的预期泄漏 = 受害者load/store地址 + 控制流变化",
    "+ CPU允许这些基本泄漏（如sysret后缓存不刷新→用户可观察内核地址）",
    "!但不允许超出此范围的更深层泄漏（如MDS泄漏store buffer数据内容）",
])

content_slide("模型实现：Unicorn模拟器中的多域模型", sub_title="§6.3", bullets=[
    "~简化策略",
    "所有Actor在用户态执行",
    "域转换模拟为简单跳转",
    "仅实现最小必要模拟（VM exit对应特权指令等）",
    "~推测行为过近似",
    "条件跳转时'临时走错误分支'模拟推测执行",
    "~合同轨迹记录",
    "所有load/store地址 + 控制流变化 = 合同允许泄漏",
    "~为什么简化模型仍然有效",
    "MRT比较的是信息暴露量而非轨迹本身",
    "简化模型只要能过滤已知（合同允许）泄漏，剩余差异就是未知泄漏",
])

content_slide("χ²统计方法：消除域转换噪声", sub_title="§7", bullets=[
    "~问题",
    "域转换（VM enter/exit）涉及大量不透明微码，引入显著噪声",
    "两条硬件轨迹可能因噪声而不同（即使输入相同）→大量假阳性",
    "~核心思路",
    "从单次轨迹比较提升到分布比较",
    "~统计工具",
    "Pearson χ²检验（将轨迹视为类别数据）",
    "χ² = Σ (obs_t(t) − N·P(t))² / (N·P(t))",
    "~自适应样本量",
    "从N=15开始 → 若检测到违反则逐步增大：40 → 160 → 320",
    "仅当违反在最大样本量仍持续才报告",
    "+检测已知泄漏时间减少一个数量级，假阳性率仅0.7%",
])

section_slide("三、执行器设计")

content_slide("执行器：从零实现定制内核模块", sub_title="§5", bullets=[
    "~为什么不用Linux现有基础设施",
    "Linux的VM/页表管理面向完整VM，开销过大",
    "Linux内核有大量安全/稳定性防护，阻止自由配置实验",
    "→ 从零实现3500+行代码的定制内核模块",
    "~两大设计需求",
    "完全可配置性：自由配置页表权限、CPU模式、VM配置，甚至允许无效系统配置",
    "低开销：目标≥100次测量/秒，需要在毫秒级创建VM和页表",
])

content_slide("执行器高层流程",
    code_block=
"disable_interrupts();            // 禁用中断\n"
"preserve_host_state();           // 保存host OS状态\n"
"configure_system_registers();    // 配置CPU（MSR、VMX/SVM等）\n"
"create_virtual_machines();       // 创建VM\n"
"create_page_tables();            // 创建页表\n"
"flush_caches_and_buffers();      // 刷新缓存\n\n"
"for (input in inputs)\n"
"    set_actor_data_memory(input);\n"
"    set_page_table_permissions();\n"
"    htrace = execute(test_case); // 执行并收集轨迹\n"
"    traces.append(htrace);\n\n"
"restore_host_state();\n"
"enable_interrupts();",
    bullets=[
    "~关键子设计",
    "状态保存/恢复：CR0、EFER、VMX控制寄存器等，防host OS崩溃",
    "故障隔离：捕获所有异常→自定义处理器；NMI→恢复后退出",
    "CPU配置：通过配置文件控制MSR、性能计数器（如禁用预取器）",
])

content_slide("执行器关键子设计", sub_title="§5.2.3",
    table=[
        ["子设计", "说明"],
        ["VM Actor配置", "启用VMX/SVM→创建VMCS/VMCB→创建EPT/NPT和guest页表→设置权限"],
        ["User Actor配置", "更新系统进入/退出点指向user代码→更新页表权限允许用户访问"],
        ["内存别名", "可选：所有VM共享相同虚拟地址布局，模拟Foreshadow-VM攻击条件"],
        ["防噪声措施", "禁用中断、禁用预取器、LFENCE防直行推测、保留寄存器子集"],
        ["目标速率", "800–4500次测量/秒"],
    ])

section_slide("四、实验结果")

content_slide("实验设计", sub_title="§8",
    table=[
        ["维度", "内容"],
        ["隔离边界", "H2V（Host→VM）、V2V（VM→VM）、K2U（Kernel→User）、U2U（User→User）"],
        ["测试目标", "Int1(KabyLake R)、Int2(Coffee Lake)、Int3(RaptorCove)、AMD1(Zen1)、AMD2(Zen2)、AMD3(Zen4)"],
        ["配置类别", "MEM（内存错误）、COMP（计算错误/DSS）、REG（特权寄存器读取/RSRR）"],
        ["实验规模", "304个campaign、88机器天、近3000万程序、超200亿硬件测量"],
    ])

content_slide("新发现①：跨VM泄漏（CVE-2024-36357）", sub_title="高安全影响", bullets=[
    "~基本信息",
    "影响CPU：AMD3 (Zen4)  |  边界：V2V  |  配置：MEM W-bit、EPT W-bit",
    "~攻击机制",
    "攻击者VM加载地址时，若缓存中包含受害VM相同虚拟地址（不同物理地址）的值",
    "受害者的缓存值影响攻击者后续操作时序 → 逐位推断受害VM内存",
    "~攻击gadget",
    "别名加载→两次写异常→两次非故障读→基于最后读取结果的缓存探测",
    "!潜在高影响：攻击者可逐位读取另一VM的任意内存",
    "+AMD在论文embargo期间已开发并部署缓解措施",
])

content_slide("新发现②：内核→用户泄漏（CVE-2024-36350）", sub_title="中等至高安全影响", bullets=[
    "~基本信息",
    "影响CPU：AMD3 (Zen4)  |  边界：K2U  |  配置：MEM U-bit",
    "~攻击机制",
    "用户进程可观察到最近N次（N≈32）内核态store操作的任意一位",
    "攻击者可通过额外store选择目标store，通过移位选择泄漏位",
    "!中高影响：攻击者可观察切换到用户态前的特权store数据",
])

content_slide("新发现③④：RSRR变体", sub_title="低安全影响但反驳AMD声明",
    two_col=(
        [
            "~CVE-2024-36349：RDTSCP-AUX",
            "影响：Int1/Int2/Int3/AMD1",
            "用户态推测性读取RDTSCP AUX寄存器",
            "即使CR4.TSD已设置（应触发#GP）",
            "!反驳AMD白皮书：声称不受RSRR影响",
            "AMD态度：不计划修复",
        ],
        [
            "~CVE-2024-36348：SMSW推测执行",
            "影响：AMD2(Zen2)/AMD3(Zen4)",
            "用户态推测性执行SMSW指令",
            "即使UMIP已启用（应触发#GP）",
            "!实质上绕过UMIP功能",
            "!同样反驳AMD不受RSRR影响的声明",
            "AMD态度：不计划修复",
        ]
    ))

content_slide("已知泄漏验证：零漏检", sub_title="§8.3",
    table=[
        ["已知攻击", "配置", "检测结果"],
        ["MDS及其变体", "MEM A/D-bit, EPT A/D-bit", "+检测成功"],
        ["Foreshadow及其RMW变体", "MEM P-bit", "+检测成功"],
        ["Meltdown", "MEM U-bit", "+检测成功"],
        ["DSS", "COMP DSS", "+检测成功"],
        ["Rogue System Register Read", "REG FSGS", "+检测成功（含变体Meltdown-3a）"],
        ["结论", "—", "+所有预期违规配置均成功检测——零漏检"],
    ])

content_slide("补丁有效性验证", sub_title="§8.4",
    table=[
        ["补丁", "结果", "说明"],
        ["VERW", "+有效", "Intel推荐，刷store buffer"],
        ["L1D_FLUSH_CMD", "+有效", "刷L1数据缓存"],
        ["WBINVD", "!无效", "非Intel推荐"],
        ["DSS当前补丁（1÷1除法）", "+有效", "sysret前假除法"],
        ["DSS旧补丁（异常处理器中除法）", "!无效", "差点合入Linux主线！"],
        ["KPTI 清除Present位", "!无效", "仍可Meltdown读取"],
        ["KPTI 完全unmapping", "+阻止Meltdown", "但仍有MDS违规"],
        ["KPTI unmapping + VERW", "+有效", "同时阻止Meltdown和MDS"],
    ])

content_slide("关键补丁发现", bullets=[
    "!DSS旧补丁无效——差点合入Linux主线内核！",
    "将除法放在#DE异常处理器而非系统进入例程中",
    "仅经安全社区审查后才被替换",
    "~这凸显了自动化验证工具的价值",
    "",
    "!KPTI仅清除Present位不能阻止Meltdown",
    "必须完全unmapping（零化PTE）才能阻止",
    "但完全unmapping仍有MDS违规 → 需配合VERW",
    "",
    "!WBINVD无效（虽功能类似缓存刷新）",
    "非Intel推荐的MDS缓解措施",
    "",
    "~核心论点验证",
    "没有自动化验证，你无法确认补丁真的有效",
])

content_slide("性能指标", sub_title="§8.5",
    table=[
        ["指标", "数据"],
        ["测量速率", "800–4500次/秒"],
        ["测试吞吐量", "60–700 test case/秒"],
        ["完整campaign耗时", "MEM: 2–24h, COMP: 3–8h, REG: 1.5–11h"],
        ["违规检测时间", "85%在1h内, 98%在4h内"],
        ["按test case计", "85%在20K轮内, 93%在40K轮内"],
        ["漏检率", "4%（10次不同seed重跑）"],
        ["假阳性", "304个campaign仅2次"],
    ])

section_slide("五、局限与未来")

content_slide("工具局限", sub_title="§8.6",
    table=[
        ["局限", "说明"],
        ["不支持Actor间共享内存", "架构级信息流（如共享内存通信）未被覆盖"],
        ["不支持自修改代码", "—"],
        ["控制结构位置固定", "IDT/GDT位置固定，无法检测CPU对控制结构位置的预测泄漏"],
        ["仅覆盖提取型攻击", "未覆盖跨域分支训练等控制流攻击（Spectre V2/LVI）"],
        ["某些泄漏需大量测量", "实验时长可能不足以发现所有可能泄漏"],
    ])

content_slide("Spectre V2覆盖分析", sub_title="FAQ Q4", bullets=[
    "~为什么无法覆盖Spectre V2",
    "攻击机制不同：提取型(观察残留) vs 注入型(污染BTB/RSB→操控受害者推测执行)",
    "模板顺序：需要attacker→victim→attacker三段式（当前为victim→attacker）",
    "非干扰性质：当前假设'攻击者数据不影响测量'，但Spectre V2恰恰利用攻击者数据",
    "模型不模拟BTB/RSB：无法预测跨域训练后受害者走哪个分支",
    "~初步方案：逆向模板",
    "攻击者训练BTB → 域转换到受害者 → 受害者执行 → 域转换回攻击者 → 测量缓存",
    "~可行性",
    "+逆向模板、统计框架、测量框架可直接复用",
    "+缓解测试(IBPB/eIBRS/RSB填充/retpoline)与现有思路一致",
    "!BTB/RSB模型精度是最大风险——建议结构化过近似",
    "预计工作量约为当前工具核心开发的30%-50%",
])

content_slide("未来方向", sub_title="§9", bullets=[
    "~短期",
    "覆盖注入型泄漏（LVI、Spectre V2）——逆向模板方案",
    "实现共享内存支持、自修改代码支持、控制结构位置随机化",
    "~中期",
    "移植到ARM、RISC-V等其他架构",
    "支持其他隔离抽象（如Intel MPK）",
    "~长期",
    "作为PoC原型工具——安全研究者可用手动模板快速构建攻击PoC",
    "Foreshadow仅需<30行模板汇编即可演示",
])

section_slide("六、总结")

content_slide("关键结论", bullets=[
    "1. 微架构隔离验证亟需自动化工具——被动式补丁验证实践屡次失败",
    "",
    "2. Actor抽象 + MRT方法提供了系统性跨域微架构泄漏检测的可行路径",
    "",
    "3. 四大新发现（CVE-2024-36357/36350/36349/36348）证明了工具有效性",
    "   其中两个高影响（跨VM泄漏、内核→用户泄漏）",
    "",
    "4. AMD的RSRR安全声明被实证反驳——在AMD CPU上发现两类RSRR漏洞",
    "",
    "5. 补丁验证价值显著——发现DSS旧补丁无效、KPTI多种变体有效性差异",
    "",
    "~范式转变",
    "从被动补丁验证 → 主动安全验证，是处理器安全设计的未来方向",
])

content_slide("研究团队与开源",
    two_col=(
        [
            "~作者",
            "Oleksii Oleksenko* (Azure Research)",
            "Flavien Solt* (ETH Zurich)",
            "Cédric Fournet (Azure Research)",
            "Jana Hofmann (MPI-SP)",
            "Boris Köpf (Azure Research)",
            "Stavros Volos (Azure Research)",
            "* 共同第一作者",
        ],
        [
            "~核心机构",
            "Microsoft Azure Research",
            "ETH Zurich",
            "MPI-SP",
            "",
            "~开源工具",
            "Revizor v1.3+",
            "github.com/microsoft/sca-fuzzer",
            "",
            "~CVE编号",
            "CVE-2024-36357 / 36350 / 36349 / 36348",
        ]
    ))

title_slide(
    "谢谢",
    "Questions & Discussion"
)

out_path = r"C:\Users\pseclab\Desktop\zhangjuchuan\github\insights\sp2026-cross-domain-test\Enter-Exit-SP26-技术演讲PPT.pptx"
prs.save(out_path)
print(f"Saved to {out_path}")
print(f"Total slides: {len(prs.slides)}")
